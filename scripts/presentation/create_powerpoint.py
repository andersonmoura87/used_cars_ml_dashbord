from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

def create_powerpoint():
    """
    Cria apresentação em PowerPoint do projeto
    """
    prs = Presentation()
    
    # Configurações de slide
    title_slide_layout = prs.slide_layouts[0]  # título
    section_slide_layout = prs.slide_layouts[1]  # título e conteúdo
    content_slide_layout = prs.slide_layouts[2]  # conteúdo
    
    # Slide 1: Capa
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Análise do Mercado de Carros Usados"
    subtitle.text = "Uma Abordagem Completa: Da Preparação dos Dados à Visualização"
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Agenda"
    content.text = """
    1. Introdução (5 min)
    2. Preparação e Limpeza dos Dados (10 min)
    3. Análise SQL e Insights Iniciais (10 min)
    4. Análise Estatística Descritiva (10 min)
    5. Visualizações e Dashboard (8 min)
    6. Demonstração do Streamlit (2 min)
    7. Perguntas e Respostas
    """
    
    # Slide 3: Introdução
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "1. Introdução"
    content.text = """
    • Contextualização do Projeto
      - Mercado de carros usados em crescimento
      - Importância da análise de dados
      - Objetivos do projeto
    
    • Visão Geral dos Dados
      - Fonte e estrutura
      - Desafios identificados
    
    • Metodologia
      - Abordagem ETL
      - Ferramentas utilizadas
    """
    
    # Slide 4: Preparação dos Dados
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "2. Preparação e Limpeza dos Dados"
    content.text = """
    • Análise Exploratória Inicial
      - Identificação de valores ausentes
      - Detecção de anomalias
    
    • Processo de Limpeza
      - Tratamento de valores nulos
      - Padronização de formatos
    
    • Criação da ABT
      - Estrutura final
      - Features criadas
    """
    
    # Slide 5: Análise SQL
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "3. Análise SQL e Insights Iniciais"
    content.text = """
    • Consultas Desenvolvidas
      - Preço médio por fabricante
      - Modelos mais anunciados
      - Análise de combustíveis
      - Análise regional
      - Análise de transmissão
    
    • Insights Principais
      - Padrões identificados
      - Distribuição geográfica
    """
    
    # Slide 6: Análise Estatística
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "4. Análise Estatística Descritiva"
    content.text = """
    • Distribuição de Preços
      - Medidas de tendência central
      - Medidas de dispersão
    
    • Análises Bivariadas
      - Correlações importantes
      - Padrões temporais
    
    • Identificação de Outliers
      - Método IQR
      - Impacto nas análises
    """
    
    # Slide 7: Dashboard
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "5. Visualizações e Dashboard"
    content.text = """
    • Power BI Dashboard
      - Estrutura e medidas DAX
      - Interatividade
    
    • Principais Visualizações
      - Box plot de preços
      - Mapa interativo
      - Top fabricantes
      - Análise de combustível
      - Dispersão preço vs. km
    """
    
    # Slide 8: Streamlit
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "6. Demonstração do Streamlit"
    content.text = """
    • Aplicação Web
      - Estrutura da aplicação
      - Funcionalidades
    
    • Diferencial do Projeto
      - Acessibilidade dos dados
      - Análises em tempo real
    """
    
    # Slide 9: Conclusões
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusões"
    content.text = """
    • Principais Descobertas
      - Padrões de preço
      - Preferências do mercado
      - Insights geográficos
    
    • Próximos Passos
      - Melhorias sugeridas
      - Novas análises
    """
    
    # Slide 10: Perguntas
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Perguntas?"
    subtitle.text = "Obrigado pela atenção!"
    
    # Criar diretório e salvar
    presentation_dir = Path("presentation")
    presentation_dir.mkdir(exist_ok=True)
    prs.save(presentation_dir / "apresentacao.pptx")
    
    print("Apresentação PowerPoint gerada com sucesso!")

if __name__ == "__main__":
    create_powerpoint() 